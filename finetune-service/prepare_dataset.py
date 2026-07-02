import json
import os
import random
import re
import time
from pathlib import Path

import google.generativeai as genai
from tqdm import tqdm

SLIDES_DIR = Path(__file__).parent / "slides"
ROOT_TRAIN_DATASET_PATH = Path(__file__).parent.parent / "train_dataset.json"
LOCAL_TRAIN_DATASET_PATH = Path(__file__).parent / "train_dataset.json"
TRAIN_JSONL_PATH = Path(__file__).parent / "train_data.jsonl"
TEST_SET_PATH = Path(__file__).parent.parent / "datasets" / "benchmark" / "test_set.json"

GEMINI_MODEL = "gemini-2.5-flash-lite"
QA_PER_SLIDE = 3
MIN_TEXT_LENGTH = 80
DELAY_BETWEEN_CALLS = 7


def normalize_text(text: str) -> str:
    if text is None:
        return ""
    lowered = text.lower().strip()
    return re.sub(r"\s+", " ", lowered)


def load_test_set_signatures() -> tuple[set[str], set[str]]:
    if not TEST_SET_PATH.exists():
        return set(), set()

    data = json.loads(TEST_SET_PATH.read_text(encoding="utf-8"))
    questions = {normalize_text(item.get("question", "")) for item in data}
    answers = {normalize_text(item.get("groundTruth", "")) for item in data}
    return questions, answers


def get_model():
    api_key = os.getenv("GEMINI_API_KEY", "").strip()
    if not api_key:
        raise RuntimeError("Thieu bien moi truong GEMINI_API_KEY")

    genai.configure(api_key=api_key)
    return genai.GenerativeModel(GEMINI_MODEL)


def extract_text_from_pptx(file_path: Path) -> list[str]:
    from pptx import Presentation

    presentation = Presentation(file_path)
    slide_texts = []

    for slide_index, slide in enumerate(presentation.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                line = " ".join(run.text for run in paragraph.runs).strip()
                if line:
                    parts.append(line)

        text = "\n".join(parts).strip()
        if len(text) >= MIN_TEXT_LENGTH:
            slide_texts.append(f"[Slide {slide_index}]\n{text}")

    return slide_texts


def extract_text_from_pdf(file_path: Path) -> list[str]:
    import PyPDF2

    page_texts = []
    with open(file_path, "rb") as file:
        reader = PyPDF2.PdfReader(file)
        for page_index, page in enumerate(reader.pages, start=1):
            text = (page.extract_text() or "").strip()
            if len(text) >= MIN_TEXT_LENGTH:
                page_texts.append(f"[Trang {page_index}]\n{text}")

    return page_texts


def extract_all_slides(slides_dir: Path) -> list[dict]:
    all_slides = []
    all_files = list(slides_dir.glob("*.pptx")) + list(slides_dir.glob("*.pdf"))

    if not all_files:
        print(f"Khong tim thay slide trong: {slides_dir}")
        return []

    print(f"Tim thay {len(all_files)} file slide")

    for file_path in all_files:
        try:
            if file_path.suffix.lower() == ".pptx":
                slide_texts = extract_text_from_pptx(file_path)
            else:
                slide_texts = extract_text_from_pdf(file_path)

            for text in slide_texts:
                all_slides.append({
                    "file": file_path.name,
                    "content": text
                })
            print(f"- {file_path.name}: {len(slide_texts)} slide/trang")
        except Exception as ex:
            print(f"- Loi doc {file_path.name}: {ex}")

    return all_slides


def generate_qa_pairs(model, slide_content: str, num_pairs: int = QA_PER_SLIDE) -> list[dict]:
    prompt = f"""Bạn là trợ lý tạo dữ liệu huấn luyện chatbot học phần.

Duới đây là nội dung slide:
---
{slide_content}
---

Hãy xin ra {num_pairs} cặp câu hỏi - trả lời dựa hoàn toàn trên nội dung slide.

Yêu cầu:
- Viết bằng tiếng Việt rõ ràng
- Câu hỏi ngắn gọn, đúng trọng tâm
- Câu trả lời đúng kiến thức trong slide
- Không thêm kiến thức bên ngoài

Trả về JSON thuần túy đúng format:
[
  {{"instruction": "Cau hoi 1", "output": "Cau tra loi 1"}},
  {{"instruction": "Cau hoi 2", "output": "Cau tra loi 2"}}
]"""

    for attempt in range(3):
        try:
            response = model.generate_content(prompt)
            raw_text = response.text.strip()
            raw_text = re.sub(r"```json\s*", "", raw_text)
            raw_text = re.sub(r"```\s*", "", raw_text).strip()

            parsed = json.loads(raw_text)
            cleaned = []

            for item in parsed:
                instruction = normalize_text(item.get("instruction", ""))
                output = normalize_text(item.get("output", ""))
                if instruction and output:
                    cleaned.append({
                        "instruction": item["instruction"].strip(),
                        "output": item["output"].strip()
                    })

            return cleaned
        except json.JSONDecodeError:
            return []
        except Exception as ex:
            if "429" in str(ex) and attempt < 2:
                time.sleep(65)
                continue
            print(f"Loi Gemini API: {ex}")
            return []

    return []


def filter_duplicate_pairs(dataset: list[dict], test_questions: set[str], test_answers: set[str]) -> tuple[list[dict], int]:
    filtered = []
    removed_count = 0

    for item in dataset:
        instruction_key = normalize_text(item.get("instruction", ""))
        answer_key = normalize_text(item.get("output", ""))

        if instruction_key in test_questions or answer_key in test_answers:
            removed_count += 1
            continue

        filtered.append({
            "instruction": item["instruction"].strip(),
            "output": item["output"].strip()
        })

    return filtered, removed_count


def save_json_dataset(dataset: list[dict]) -> None:
    ROOT_TRAIN_DATASET_PATH.write_text(
        json.dumps(dataset, ensure_ascii=False, indent=2),
        encoding="utf-8"
    )
    LOCAL_TRAIN_DATASET_PATH.write_text(
        json.dumps(dataset, ensure_ascii=False, indent=2),
        encoding="utf-8"
    )


def save_jsonl_dataset(dataset: list[dict]) -> None:
    lines = []
    for item in dataset:
        lines.append(json.dumps({
            "instruction": item["instruction"],
            "input": "",
            "output": item["output"]
        }, ensure_ascii=False))

    TRAIN_JSONL_PATH.write_text("\n".join(lines), encoding="utf-8")


def verify_output(dataset: list[dict]) -> None:
    if not dataset:
        print("Khong co du lieu de verify")
        return

    for index, item in enumerate(dataset):
        if not item.get("instruction", "").strip():
            raise RuntimeError(f"Item {index} thieu instruction")
        if not item.get("output", "").strip():
            raise RuntimeError(f"Item {index} thieu output")

    print(f"JSON hop le: {len(dataset)} cap")
    print(f"Da ghi: {ROOT_TRAIN_DATASET_PATH}")
    print(f"Da ghi: {LOCAL_TRAIN_DATASET_PATH}")
    print(f"Da ghi: {TRAIN_JSONL_PATH}")


def run_pipeline():
    print("BAT DAU CHUAN BI DATASET FINE-TUNE")
    SLIDES_DIR.mkdir(exist_ok=True)

    test_questions, test_answers = load_test_set_signatures()
    slides = extract_all_slides(SLIDES_DIR)
    if not slides:
        print("Khong co slide de xu ly")
        return

    model = get_model()
    raw_dataset = []

    for slide in tqdm(slides, desc="Dang xu ly slide"):
        pairs = generate_qa_pairs(model, slide["content"], QA_PER_SLIDE)
        raw_dataset.extend(pairs)
        time.sleep(DELAY_BETWEEN_CALLS)

    random.shuffle(raw_dataset)
    filtered_dataset, removed_count = filter_duplicate_pairs(raw_dataset, test_questions, test_answers)

    save_json_dataset(filtered_dataset)
    save_jsonl_dataset(filtered_dataset)
    verify_output(filtered_dataset)

    print(f"So slide/trang da doc: {len(slides)}")
    print(f"So cap QA sinh ra: {len(raw_dataset)}")
    print(f"So cap bi loai do trung test set: {removed_count}")
    print(f"So cap cuoi cung: {len(filtered_dataset)}")


if __name__ == "__main__":
    run_pipeline()
